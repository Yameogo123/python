import os
import pptx
import unicodedata

#obtenir le chemin dacces vers cv
cv=os.chdir("cv")
#recuprer la liste des cv disponible
listofCv=os.listdir(cv)


#methode pour lire fichier pptx
def getPptx(file):
    presentation = pptx.Presentation(file)
    results = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text!="":
                        results.append(paragraph.text)
    return results

#Permet de recuperer le cv en fonction de la compétence
def getCV(Task):
    l=[]
    liste = []
    for ask in Task:
        s2 = unicodedata.normalize('NFD', ask).encode('ascii', 'ignore').decode("utf-8")
        for cv in listofCv:
            #obtenir l'extension
            extension=os.path.splitext(cv)[1]
            #comparer les extensions
            if extension == '.ppt' or extension == '.pptx':
                file=getPptx(cv)
                for f in file:
                    if s2 in f.lower() or ask in f.lower():
                        liste.append(cv)
    for cv in liste:
        if liste.count(cv)==len(Task) and cv not in l:
            l.append(cv)
    return l

#recupere l'adresse de chaque CV de la compétence
def getAdress(Task):
    listes=getCV(Task)
    addresse=[]
    for cv in listes:
        output=[cv]
        file = getPptx(cv)
        for f in file:
            if "adresse" in f.lower() or "addresse" in f.lower() or "rue" in f.lower():
                index=file.index(f)
                recup=file[index]+" "+file[index+1]
                output.append(recup)
                addresse.append(output)
    return addresse

#classe les cv en deux les proches et les loins
def classerCV(Task):
    villes=["dakar","ziguinchor","saint louis","touba","thies","kaolack","rufisque","pikine"]
    adresses=getAdress(Task)
    proche = []
    other=[]
    for adresse in adresses:
        if any(item in adresse[1].lower() for item in villes):
            proche.append(adresse[0])
        else:
            other.append(adresse[0])
    return proche,other


#print(classerCV(["photoshop"]))

def extract(filepath):
    prs = pptx.Presentation(filepath) 

    # Traverse
    for slide in prs.slides:
        for obj in slide.shapes:
            
            try:
                # recupère les images sous format code binaire 
                imdata = obj.image.blob
                # récupère le type de l'image et son extension
                imagetype = obj.image.content_type
                typekey = imagetype.find('/') + 1
                imtype = imagetype[typekey:]
                # stock les image de chaque cv dans un dossier portant son nom
                path = "../"+filepath[:filepath.find(".")+1]+"/"
                if not os.path.exists(path):
                    os.makedirs(path)
                image_file = path + obj.name + "." + imtype
                file_str = open(image_file, 'wb')
                file_str.write(imdata)
                file_str.close()

            except:
                pass

def extractAll():
    for cv in listofCv:
        extract(cv)

extractAll()